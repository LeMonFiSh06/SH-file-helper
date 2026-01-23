from __future__ import annotations

import enum
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pdf2docx import Converter
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


class ConversionMode(str, enum.Enum):
    PPTX_TO_PDF = "pptx_to_pdf"
    PDF_TO_PPTX = "pdf_to_pptx"
    DOCX_TO_PDF = "docx_to_pdf"
    PDF_TO_DOCX = "pdf_to_docx"
    IMAGE_TO_PDF = "image_to_pdf"
    IMAGES_TO_PDF = "images_to_pdf"


@dataclass(frozen=True)
class ConversionRequest:
    mode: ConversionMode
    input_paths: Sequence[Path]
    output_path: Path


class ConversionError(RuntimeError):
    pass


def convert(request: ConversionRequest) -> Path:
    """Run a conversion based on the selected mode.

    The conversion relies on LibreOffice for Office <-> PDF transforms and
    Pillow for image-to-PDF conversions.
    """
    if request.mode in {
        ConversionMode.PPTX_TO_PDF,
        ConversionMode.DOCX_TO_PDF,
    }:
        _convert_with_libreoffice(request)
    elif request.mode == ConversionMode.PDF_TO_DOCX:
        _convert_pdf_to_docx(request)
    elif request.mode == ConversionMode.PDF_TO_PPTX:
        _convert_pdf_to_pptx(request)
    elif request.mode == ConversionMode.IMAGE_TO_PDF:
        _convert_single_image_to_pdf(request.input_paths, request.output_path)
    elif request.mode == ConversionMode.IMAGES_TO_PDF:
        _convert_multiple_images_to_pdf(request.input_paths, request.output_path)
    else:
        raise ConversionError(f"Unsupported conversion mode: {request.mode}")

    return request.output_path


def _convert_with_libreoffice(request: ConversionRequest) -> None:
    if len(request.input_paths) != 1:
        raise ConversionError("LibreOffice conversions require exactly one input file.")

    input_path = request.input_paths[0]
    output_dir = request.output_path.parent
    output_dir.mkdir(parents=True, exist_ok=True)

    target_format = _libreoffice_target_format(request.mode)
    command = [
        _libreoffice_executable(),
        "--headless",
        "--convert-to",
        target_format,
        "--outdir",
        str(output_dir),
        str(input_path),
    ]

    try:
        subprocess.run(command, check=True, capture_output=True, text=True)
    except FileNotFoundError as exc:
        raise ConversionError(
            "LibreOffice not found. Install LibreOffice and ensure 'soffice' is on PATH."
        ) from exc
    except subprocess.CalledProcessError as exc:
        raise ConversionError(
            f"LibreOffice conversion failed: {exc.stderr or exc.stdout}"
        ) from exc

    produced = output_dir / input_path.with_suffix(f".{target_format}").name
    if not produced.exists():
        raise ConversionError("LibreOffice did not create the expected output file.")

    if produced != request.output_path:
        produced.replace(request.output_path)


def _libreoffice_target_format(mode: ConversionMode) -> str:
    mapping = {
        ConversionMode.PPTX_TO_PDF: "pdf",
        ConversionMode.DOCX_TO_PDF: "pdf",
    }
    try:
        return mapping[mode]
    except KeyError as exc:
        raise ConversionError(f"Unsupported LibreOffice conversion mode: {mode}") from exc


def _libreoffice_executable() -> str:
    return "soffice"


def _convert_single_image_to_pdf(input_paths: Sequence[Path], output_path: Path) -> None:
    if len(input_paths) != 1:
        raise ConversionError("Image-to-PDF conversion requires exactly one image.")

    image_path = input_paths[0]
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with Image.open(image_path) as image:
        rgb = image.convert("RGB")
        rgb.save(output_path, format="PDF")


def _convert_multiple_images_to_pdf(input_paths: Iterable[Path], output_path: Path) -> None:
    image_paths = list(input_paths)
    if len(image_paths) < 2:
        raise ConversionError("Multi-image PDF conversion requires at least two images.")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    images = []
    for path in image_paths:
        with Image.open(path) as image:
            images.append(image.convert("RGB"))

    cover, *rest = images
    cover.save(output_path, format="PDF", save_all=True, append_images=rest)


def _convert_pdf_to_docx(request: ConversionRequest) -> None:
    if len(request.input_paths) != 1:
        raise ConversionError("PDF-to-DOCX conversion requires exactly one input file.")

    input_path = request.input_paths[0]
    output_path = request.output_path
    output_path.parent.mkdir(parents=True, exist_ok=True)

    converter = Converter(str(input_path))
    try:
        converter.convert(str(output_path), start=0)
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(f"PDF-to-DOCX conversion failed: {exc}") from exc
    finally:
        converter.close()


def _convert_pdf_to_pptx(request: ConversionRequest) -> None:
    if len(request.input_paths) != 1:
        raise ConversionError("PDF-to-PPTX conversion requires exactly one input file.")

    input_path = request.input_paths[0]
    output_path = request.output_path
    output_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        images = convert_from_path(str(input_path), dpi=200)
    except Exception as exc:  # noqa: BLE001
        raise ConversionError(
            "PDF-to-PPTX conversion failed while rendering pages. Ensure Poppler is installed."
        ) from exc

    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    for image in images:
        slide = presentation.slides.add_slide(blank_layout)
        width, height = image.size
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        image_stream = _image_to_stream(image)
        slide.shapes.add_picture(image_stream, Inches(0), Inches(0), slide_width, slide_height)
        image_stream.close()

    presentation.save(str(output_path))


def _image_to_stream(image: Image.Image):
    from io import BytesIO

    stream = BytesIO()
    image.convert("RGB").save(stream, format="PNG")
    stream.seek(0)
    return stream