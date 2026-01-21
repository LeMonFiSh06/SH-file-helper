from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List

from docx import Document
from pptx import Presentation

from ocr import OcrRequest, ocr_pdf


@dataclass(frozen=True)
class TextExtractRequest:
    input_path: Path
    language: str = "eng"
    dpi: int = 300


class TextExtractError(RuntimeError):
    pass


def extract_text(request: TextExtractRequest) -> str:
    if not request.input_path.exists():
        raise TextExtractError(f"Input file not found: {request.input_path}")

    suffix = request.input_path.suffix.lower()
    if suffix == ".txt":
        return request.input_path.read_text(encoding="utf-8")
    if suffix == ".docx":
        return _extract_docx_text(request.input_path)
    if suffix == ".pptx":
        return _extract_pptx_text(request.input_path)
    if suffix == ".pdf":
        return _extract_pdf_text(request)

    raise TextExtractError(f"Unsupported input type for glossary: {suffix}")


def _extract_docx_text(path: Path) -> str:
    try:
        document = Document(str(path))
    except Exception as exc:  # noqa: BLE001
        raise TextExtractError(f"Failed to read DOCX: {path}") from exc
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)


def _extract_pptx_text(path: Path) -> str:
    try:
        # Pylance 兼容：显式转 str
        presentation = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        raise TextExtractError(f"Failed to read PPTX: {path}") from exc

    chunks: List[str] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame  # type: ignore
                if tf is None:
                    continue

                for para in tf.paragraphs:
                    text = (para.text or "").strip()
                    if text:
                        chunks.append(text)

    return "\n".join(chunks)



def _extract_pdf_text(request: TextExtractRequest) -> str:
    result = ocr_pdf(OcrRequest(request.input_path, request.language, request.dpi))
    return result.text
