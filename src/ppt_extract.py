from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Iterable, List

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract


@dataclass(frozen=True)
class PptExtractRequest:
    input_path: Path
    language: str = "eng"


@dataclass(frozen=True)
class SlideText:
    slide_number: int
    lines: List[str]


@dataclass(frozen=True)
class PptExtractResult:
    slides: List[SlideText]

    def to_text(self) -> str:
        sections = []
        for slide in self.slides:
            sections.append(f"Slide {slide.slide_number}")
            sections.extend(slide.lines)
        return "\n".join(sections)


class PptExtractError(RuntimeError):
    pass


def extract_ppt_text(request: PptExtractRequest) -> PptExtractResult:
    if not request.input_path.exists():
        raise PptExtractError(f"Input file not found: {request.input_path}")

    try:
        presentation = Presentation(str(request.input_path))
    except Exception as exc:  # noqa: BLE001
        raise PptExtractError(f"Failed to read PPTX: {request.input_path}") from exc

    slides: List[SlideText] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: List[str] = []
        for shape in slide.shapes:
            lines.extend(_extract_shape_text(shape, request.language))
        slides.append(SlideText(slide_number=index, lines=_dedupe_lines(lines)))

    return PptExtractResult(slides=slides)


def _extract_shape_text(shape, language: str) -> List[str]:
    lines: List[str] = []
    if hasattr(shape, "text"):
        for line in shape.text.splitlines():
            cleaned = line.strip()
            if cleaned:
                lines.append(cleaned)

    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    lines.append(cell_text)

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = Image.open(BytesIO(shape.image.blob))
        except Exception:  # noqa: BLE001
            return lines
        ocr_text = pytesseract.image_to_string(image, lang=language)
        for line in ocr_text.splitlines():
            cleaned = line.strip()
            if cleaned:
                lines.append(cleaned)

    return lines


def _dedupe_lines(lines: Iterable[str]) -> List[str]:
    seen = set()
    output: List[str] = []
    for line in lines:
        if line in seen:
            continue
        seen.add(line)
        output.append(line)
    return output
